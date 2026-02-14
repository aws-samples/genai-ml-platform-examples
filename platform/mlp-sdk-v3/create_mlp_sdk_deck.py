#!/usr/bin/env python3
"""
Create ML Platform SDK presentation using ppt_template.pptx
Analyzes template structure and applies appropriate layouts
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def analyze_template(template_path):
    """Analyze template to understand available layouts"""
    prs = Presentation(template_path)
    
    print(f"\n{'='*60}")
    print(f"Template Analysis: {template_path}")
    print(f"{'='*60}\n")
    
    print(f"Total slide layouts: {len(prs.slide_layouts)}\n")
    
    for idx, layout in enumerate(prs.slide_layouts):
        print(f"Layout {idx}: {layout.name}")
        print(f"  Placeholders: {len(layout.placeholders)}")
        for placeholder in layout.placeholders:
            print(f"    - {placeholder.placeholder_format.idx}: {placeholder.name}")
        print()
    
    return prs

def add_title_slide(prs, title, subtitle):
    """Add title slide using template layout"""
    # Use layout 0 for title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = title
    
    # Set subtitle if placeholder exists
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    
    return slide

def add_agenda_slide(prs, title, agenda_items):
    """Add agenda slide using template layout"""
    # Try to find agenda layout, otherwise use content layout
    slide_layout = None
    for layout in prs.slide_layouts:
        if 'agenda' in layout.name.lower():
            slide_layout = layout
            break
    
    if not slide_layout:
        slide_layout = prs.slide_layouts[1]  # Use content layout
    
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = title
    
    # Add agenda items to content placeholder
    for placeholder in slide.placeholders:
        if placeholder.placeholder_format.type == 2:  # Body placeholder
            tf = placeholder.text_frame
            tf.clear()
            
            for i, item in enumerate(agenda_items, 1):
                p = tf.add_paragraph()
                p.text = f"{i}. {item}"
                p.level = 0
                if p.font:
                    p.font.size = Pt(20)
            break
    
    return slide

def add_content_slide(prs, title, content_items):
    """Add content slide with bullets using template layout"""
    # Use layout 1 for content
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = title
    
    # Add content to body placeholder
    for placeholder in slide.placeholders:
        if placeholder.placeholder_format.type == 2:  # Body placeholder
            tf = placeholder.text_frame
            tf.clear()
            
            for item in content_items:
                if isinstance(item, dict):
                    p = tf.add_paragraph()
                    p.text = item['text']
                    p.level = item.get('level', 0)
                    if p.font:
                        p.font.size = Pt(item.get('font_size', 18))
                elif item == "":
                    # Empty line for spacing
                    p = tf.add_paragraph()
                    p.text = ""
                else:
                    p = tf.add_paragraph()
                    p.text = item
                    p.level = 0
                    if p.font:
                        p.font.size = Pt(18)
            break
    
    return slide

def add_two_column_slide(prs, title, left_content, right_content):
    """Add two-column content slide"""
    # Try to find two-column layout
    slide_layout = None
    for layout in prs.slide_layouts:
        if 'two' in layout.name.lower() or 'comparison' in layout.name.lower():
            slide_layout = layout
            break
    
    if not slide_layout:
        slide_layout = prs.slide_layouts[1]  # Fallback to content layout
    
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = title
    
    # If template has two content placeholders, use them
    content_placeholders = [p for p in slide.placeholders if p.placeholder_format.type == 2]
    
    if len(content_placeholders) >= 2:
        # Use template placeholders
        for i, (placeholder, content) in enumerate(zip(content_placeholders[:2], [left_content, right_content])):
            tf = placeholder.text_frame
            tf.clear()
            for item in content:
                p = tf.add_paragraph()
                p.text = item
                if p.font:
                    p.font.size = Pt(16)
    else:
        # Create custom text boxes
        left = Inches(0.5)
        top = Inches(2)
        width = Inches(4.25)
        height = Inches(4.5)
        
        # Left column
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        for item in left_content:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(16)
        
        # Right column
        left = Inches(5.25)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        for item in right_content:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(16)
    
    return slide

def create_mlp_sdk_presentation():
    """Create the ML Platform SDK presentation"""
    
    template_path = "ppt_template.pptx"
    
    # Analyze template first
    print("Analyzing template structure...")
    prs = analyze_template(template_path)
    
    print("\nCreating ML Platform SDK presentation...")
    print("="*60)
    
    # Slide 1: Title
    add_title_slide(
        prs,
        "Building an ML Platform SDK Wrapper",
        "Simplifying SageMaker Operations with Configuration-Driven Infrastructure"
    )
    print("✓ Slide 1: Title")
    
    # Slide 2: Agenda
    add_agenda_slide(
        prs,
        "Agenda",
        [
            "What: Understanding the Challenge",
            "Why: The Need for Simplification",
            "How: ML Platform SDK Architecture",
            "Configuration & Setup",
            "Core Features & Examples",
            "Advanced Capabilities",
            "Getting Started"
        ]
    )
    print("✓ Slide 2: Agenda")
    
    # Slide 3: What - The Challenge
    add_content_slide(
        prs,
        "What: The ML Infrastructure Challenge",
        [
            "Data scientists face significant infrastructure overhead:",
            "",
            "Complex Configuration Requirements",
            {"text": "VPC settings, security groups, IAM roles", "level": 1},
            {"text": "S3 bucket policies, KMS encryption keys", "level": 1},
            {"text": "Network isolation and compliance settings", "level": 1},
            "",
            "Repetitive Boilerplate Code",
            {"text": "50+ lines of infrastructure code per training job", "level": 1},
            {"text": "Copy-paste errors across projects", "level": 1},
            {"text": "Inconsistent patterns across teams", "level": 1},
            "",
            "Governance Challenges",
            {"text": "Difficult to enforce security policies", "level": 1},
            {"text": "Limited audit trails and cost attribution", "level": 1}
        ]
    )
    print("✓ Slide 3: The Challenge")
    
    # Slide 4: What - Current State
    add_content_slide(
        prs,
        "What: Current State Without SDK Wrapper",
        [
            "Every SageMaker operation requires extensive setup:",
            "",
            "Training Job Example (Traditional Approach)",
            {"text": "Import 10+ SageMaker classes", "level": 1},
            {"text": "Define Compute, Networking, OutputDataConfig", "level": 1},
            {"text": "Configure ResourceConfig, StoppingCondition", "level": 1},
            {"text": "Specify VPC, subnets, security groups", "level": 1},
            {"text": "Set IAM roles and KMS encryption", "level": 1},
            {"text": "Create ModelTrainer with all parameters", "level": 1},
            "",
            "Result: 50+ lines of boilerplate",
            "Problem: Repeated across every project",
            "Impact: Slow iteration, high error rate"
        ]
    )
    print("✓ Slide 4: Current State")
    
    # Slide 5: Why - The Need for Change
    add_content_slide(
        prs,
        "Why: The Need for Simplification",
        [
            "Three Critical Problems to Solve:",
            "",
            "1. Infrastructure Abstraction Complexity",
            {"text": "Data scientists shouldn't manage VPCs and security groups", "level": 1},
            {"text": "Focus should be on ML, not infrastructure", "level": 1},
            "",
            "2. Consistency and Standardization",
            {"text": "Different teams use different patterns", "level": 1},
            {"text": "Configuration drift between environments", "level": 1},
            {"text": "Knowledge silos and tribal knowledge", "level": 1},
            "",
            "3. Governance and Compliance",
            {"text": "Central policy enforcement needed", "level": 1},
            {"text": "Audit trails for compliance", "level": 1},
            {"text": "Cost management and attribution", "level": 1}
        ]
    )
    print("✓ Slide 5: Why Change")
    
    # Slide 6: Why - Business Impact
    add_content_slide(
        prs,
        "Why: Business Impact",
        [
            "Quantifiable Benefits of SDK Wrapper:",
            "",
            "Development Velocity",
            {"text": "90% reduction in boilerplate code", "level": 1},
            {"text": "10x faster iteration cycles", "level": 1},
            {"text": "Hours to production instead of days", "level": 1},
            "",
            "Team Productivity",
            {"text": "New team members productive in minutes", "level": 1},
            {"text": "Reduced context switching", "level": 1},
            {"text": "Focus on model quality, not infrastructure", "level": 1},
            "",
            "Risk Reduction",
            {"text": "Fewer configuration errors", "level": 1},
            {"text": "Consistent security policies", "level": 1},
            {"text": "Complete audit trails", "level": 1}
        ]
    )
    print("✓ Slide 6: Business Impact")
    
    # Slide 7: How - Solution Overview
    add_content_slide(
        prs,
        "How: ML Platform SDK (mlp_sdk)",
        [
            "Configuration-driven wrapper around SageMaker Python SDK v3",
            "",
            "Core Concept",
            {"text": "Define infrastructure once in YAML configuration", "level": 1},
            {"text": "Simple, declarative API for ML operations", "level": 1},
            {"text": "Runtime flexibility to override any setting", "level": 1},
            "",
            "Key Features",
            {"text": "Session-based interface with smart defaults", "level": 1},
            {"text": "Built-in governance with audit trails", "level": 1},
            {"text": "AES-256-GCM encryption for sensitive values", "level": 1},
            {"text": "Zero lock-in - access underlying SDK anytime", "level": 1},
            "",
            "Result: 5 lines of code instead of 50+"
        ]
    )
    print("✓ Slide 7: Solution Overview")
    
    # Slide 8: How - Architecture
    add_content_slide(
        prs,
        "How: Three-Layer Architecture",
        [
            "Intelligent wrapper architecture:",
            "",
            "Layer 1: MLP_Session",
            {"text": "Central orchestrator and configuration manager", "level": 1},
            {"text": "Loads YAML configuration", "level": 1},
            {"text": "Manages boto3 and SageMaker sessions", "level": 1},
            "",
            "Layer 2: Specialized Wrappers",
            {"text": "Training Wrapper - Model training operations", "level": 1},
            {"text": "Processing Wrapper - Data preprocessing", "level": 1},
            {"text": "Feature Store Wrapper - Feature management", "level": 1},
            {"text": "Pipeline Wrapper - ML workflow orchestration", "level": 1},
            "",
            "Layer 3: SageMaker SDK v3",
            {"text": "Full access to underlying SDK", "level": 1},
            {"text": "No functionality locked away", "level": 1}
        ]
    )
    print("✓ Slide 8: Architecture")
    
    # Slide 9: How - Configuration Precedence
    add_content_slide(
        prs,
        "How: Configuration Precedence System",
        [
            "Three-tier precedence for maximum flexibility:",
            "",
            "Priority 1: Runtime Parameters (Highest)",
            {"text": "Explicitly passed to wrapper methods", "level": 1},
            {"text": "Override everything for special cases", "level": 1},
            {"text": "Example: instance_type='ml.p3.8xlarge'", "level": 1},
            "",
            "Priority 2: YAML Configuration (Middle)",
            {"text": "Defined in admin-config.yaml", "level": 1},
            {"text": "Team-wide defaults and standards", "level": 1},
            {"text": "Environment-specific settings", "level": 1},
            "",
            "Priority 3: SageMaker SDK Defaults (Lowest)",
            {"text": "Built-in SageMaker defaults", "level": 1},
            {"text": "Fallback when nothing specified", "level": 1}
        ]
    )
    print("✓ Slide 9: Configuration Precedence")
    
    # Slide 10: Configuration Example
    add_content_slide(
        prs,
        "Configuration File Structure",
        [
            "Define infrastructure once in admin-config.yaml:",
            "",
            "S3 Configuration",
            {"text": "default_bucket, input_prefix, output_prefix", "level": 1},
            "",
            "Networking Configuration",
            {"text": "vpc_id, security_group_ids, subnets", "level": 1},
            "",
            "Compute Configuration",
            {"text": "processing_instance_type, training_instance_type", "level": 1},
            {"text": "instance_count, volume_size", "level": 1},
            "",
            "IAM Configuration",
            {"text": "execution_role ARN", "level": 1},
            "",
            "KMS Configuration (Optional)",
            {"text": "key_id for encryption", "level": 1}
        ]
    )
    print("✓ Slide 10: Configuration")
    
    # Slide 11: Code Comparison
    add_two_column_slide(
        prs,
        "Before vs After: Training Job",
        [
            "❌ Without mlp_sdk (50+ lines)",
            "",
            "• Import multiple classes",
            "• Define Compute config",
            "• Define Networking config",
            "• Define OutputDataConfig",
            "• Define ResourceConfig",
            "• Define StoppingCondition",
            "• Create ModelTrainer",
            "• Configure all parameters",
            "• Start training",
            "",
            "Verbose, error-prone, repetitive"
        ],
        [
            "✅ With mlp_sdk (5 lines)",
            "",
            "trainer = session.run_training_job(",
            "    job_name='my-training',",
            "    training_image=container,",
            "    inputs={'train': 's3://...'}",
            ")",
            "",
            "",
            "",
            "",
            "Clean, focused, maintainable"
        ]
    )
    print("✓ Slide 11: Code Comparison")
    
    # Slide 12: Core Features - Training
    add_content_slide(
        prs,
        "Core Feature: Training Jobs",
        [
            "Simplified training job execution:",
            "",
            "Basic Usage",
            {"text": "session = MLP_Session()", "level": 1},
            {"text": "trainer = session.run_training_job(...)", "level": 1},
            "",
            "What's Handled Automatically",
            {"text": "✓ Instance type from configuration", "level": 1},
            {"text": "✓ VPC and security group setup", "level": 1},
            {"text": "✓ IAM role assignment", "level": 1},
            {"text": "✓ KMS encryption configuration", "level": 1},
            {"text": "✓ Output path management", "level": 1},
            "",
            "Runtime Flexibility",
            {"text": "Override any setting when needed", "level": 1},
            {"text": "Example: Use GPU for specific job", "level": 1}
        ]
    )
    print("✓ Slide 12: Training Jobs")
    
    # Slide 13: Core Features - Processing & Feature Store
    add_content_slide(
        prs,
        "Core Features: Processing & Feature Store",
        [
            "Processing Jobs",
            {"text": "session.run_processing_job()", "level": 1},
            {"text": "Simplified input/output data mapping", "level": 1},
            {"text": "Automatic resource configuration", "level": 1},
            "",
            "Feature Store Operations",
            {"text": "session.create_feature_group()", "level": 1},
            {"text": "Validates feature definitions", "level": 1},
            {"text": "Handles online/offline store config", "level": 1},
            {"text": "Manages feature group metadata", "level": 1},
            "",
            "Model Deployment",
            {"text": "session.deploy_model()", "level": 1},
            {"text": "Automatic endpoint configuration", "level": 1},
            {"text": "VPC and security settings applied", "level": 1}
        ]
    )
    print("✓ Slide 13: Processing & Feature Store")
    
    # Slide 14: Advanced Features
    add_content_slide(
        prs,
        "Advanced Capabilities",
        [
            "Audit Trails",
            {"text": "Track all operations for debugging and compliance", "level": 1},
            {"text": "Export to JSON or CSV for analysis", "level": 1},
            {"text": "Filter by operation type or status", "level": 1},
            "",
            "Encryption Support",
            {"text": "AES-256-GCM for sensitive configuration values", "level": 1},
            {"text": "Integration with AWS KMS", "level": 1},
            {"text": "Environment variables or file-based keys", "level": 1},
            "",
            "Multi-Environment Support",
            {"text": "Separate configs for dev/staging/prod", "level": 1},
            {"text": "Environment-specific settings", "level": 1},
            "",
            "Full SDK Access",
            {"text": "No lock-in - access SageMaker SDK anytime", "level": 1},
            {"text": "Use low-level APIs when needed", "level": 1}
        ]
    )
    print("✓ Slide 14: Advanced Features")
    
    # Slide 15: Getting Started
    add_content_slide(
        prs,
        "Getting Started (5 Minutes)",
        [
            "Quick start in three steps:",
            "",
            "Step 1: Install",
            {"text": "pip install mlp_sdk", "level": 1},
            "",
            "Step 2: Generate Configuration",
            {"text": "python examples/generate_admin_config.py --interactive", "level": 1},
            {"text": "Answer prompts to create admin-config.yaml", "level": 1},
            "",
            "Step 3: Start Building",
            {"text": "from mlp_sdk import MLP_Session", "level": 1},
            {"text": "session = MLP_Session()", "level": 1},
            {"text": "trainer = session.run_training_job(...)", "level": 1},
            "",
            "Resources",
            {"text": "Examples: examples/xgboost_training_example.ipynb", "level": 1},
            {"text": "Documentation: README.md", "level": 1}
        ]
    )
    print("✓ Slide 15: Getting Started")
    
    # Save presentation
    output_file = "ML_Platform_SDK_Deck.pptx"
    prs.save(output_file)
    
    print("="*60)
    print(f"\n✅ Presentation created: {output_file}")
    print(f"   Total slides: {len(prs.slides)}")
    print(f"   Template used: {template_path}")
    
    return output_file

if __name__ == "__main__":
    try:
        output_file = create_mlp_sdk_presentation()
        print(f"\n🎉 Success! Open {output_file} to view the presentation.")
    except Exception as e:
        print(f"\n❌ Error creating presentation: {e}")
        import traceback
        traceback.print_exc()
