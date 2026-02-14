#!/usr/bin/env python3
"""
Create ML Platform SDK presentation with concise bullet points
that fit within slide boundaries
"""

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

def add_bullet_slide(prs, layout_idx, title, bullets):
    """Add a slide with title and bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    
    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = title
    
    # Find content placeholder and add bullets
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            tf = shape.text_frame
            tf.clear()
            
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
                else:
                    p = tf.add_paragraph()
                
                if isinstance(bullet, dict):
                    p.text = bullet['text']
                    p.level = bullet.get('level', 0)
                    if 'font_size' in bullet:
                        p.font.size = Pt(bullet['font_size'])
                else:
                    p.text = bullet
                    p.level = 0
            
            return slide
    
    return slide

def add_two_column_slide(prs, layout_idx, title, left_content, right_content):
    """Add a two-column comparison slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    
    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = title
    
    # Find content placeholders
    content_shapes = [s for s in slide.shapes if s.has_text_frame and s != slide.shapes.title]
    
    if len(content_shapes) >= 2:
        # Left column
        tf_left = content_shapes[0].text_frame
        tf_left.clear()
        for i, item in enumerate(left_content):
            p = tf_left.paragraphs[0] if i == 0 and tf_left.paragraphs else tf_left.add_paragraph()
            p.text = item
            if p.font:
                p.font.size = Pt(14)
        
        # Right column
        tf_right = content_shapes[1].text_frame
        tf_right.clear()
        for i, item in enumerate(right_content):
            p = tf_right.paragraphs[0] if i == 0 and tf_right.paragraphs else tf_right.add_paragraph()
            p.text = item
            if p.font:
                p.font.size = Pt(14)
    
    return slide

def create_presentation():
    """Create the ML Platform SDK presentation"""
    
    # Load template
    prs = Presentation("ppt_template.pptx")
    
    print(f"Template loaded: {len(prs.slide_layouts)} layouts available")
    
    # Define layout indices
    TITLE_LAYOUT = 1      # Title Slide 1B
    AGENDA_LAYOUT = 6     # Agenda Slide 1
    CONTENT_LAYOUT = 16   # Title and Bulleted Content
    TWO_COL_LAYOUT = 21   # Two Content
    
    print(f"\nUsing layouts:")
    print(f"  Title: {TITLE_LAYOUT} - {prs.slide_layouts[TITLE_LAYOUT].name}")
    print(f"  Agenda: {AGENDA_LAYOUT} - {prs.slide_layouts[AGENDA_LAYOUT].name}")
    print(f"  Content: {CONTENT_LAYOUT} - {prs.slide_layouts[CONTENT_LAYOUT].name}")
    print(f"  Two Column: {TWO_COL_LAYOUT} - {prs.slide_layouts[TWO_COL_LAYOUT].name}")
    print()
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_LAYOUT])
    if slide.shapes.title:
        slide.shapes.title.text = "Building an ML Platform SDK Wrapper"
    for shape in slide.placeholders:
        if shape.placeholder_format.type == 3:  # Subtitle
            shape.text = "Simplifying SageMaker with Configuration-Driven Infrastructure"
            break
    print("✓ Slide 1: Title")
    
    # Slide 2: Agenda
    add_bullet_slide(prs, AGENDA_LAYOUT, "Agenda", [
        "What: The ML Infrastructure Challenge",
        "Why: The Need for Simplification",
        "How: ML Platform SDK Architecture",
        "Configuration & Examples",
        "Getting Started"
    ])
    print("✓ Slide 2: Agenda")
    
    # Slide 3: What - The Challenge
    add_bullet_slide(prs, CONTENT_LAYOUT, "What: The ML Infrastructure Challenge", [
        "Data scientists face infrastructure overhead:",
        "",
        "Complex Configuration",
        {"text": "VPC, security groups, IAM roles", "level": 1},
        {"text": "S3 policies, KMS encryption", "level": 1},
        "",
        "Repetitive Code",
        {"text": "50+ lines per training job", "level": 1},
        {"text": "Copy-paste errors", "level": 1},
        "",
        "Governance Gaps",
        {"text": "Hard to enforce policies", "level": 1},
        {"text": "Limited audit trails", "level": 1}
    ])
    print("✓ Slide 3: The Challenge")
    
    # Slide 4: What - Current State
    add_bullet_slide(prs, CONTENT_LAYOUT, "What: Current State Without SDK", [
        "Traditional SageMaker requires extensive setup:",
        "",
        "Training Job Example",
        {"text": "Import 10+ classes", "level": 1},
        {"text": "Configure Compute, Networking, Output", "level": 1},
        {"text": "Set VPC, subnets, security groups", "level": 1},
        {"text": "Specify IAM roles and KMS keys", "level": 1},
        "",
        "Result: 50+ lines of boilerplate",
        "Problem: Repeated across projects",
        "Impact: Slow iteration, high errors"
    ])
    print("✓ Slide 4: Current State")
    
    # Slide 5: Why - The Need
    add_bullet_slide(prs, CONTENT_LAYOUT, "Why: The Need for Simplification", [
        "Three critical problems:",
        "",
        "1. Infrastructure Complexity",
        {"text": "Focus on ML, not infrastructure", "level": 1},
        "",
        "2. Consistency",
        {"text": "Different teams, different patterns", "level": 1},
        {"text": "Configuration drift", "level": 1},
        "",
        "3. Governance",
        {"text": "Central policy enforcement", "level": 1},
        {"text": "Audit trails for compliance", "level": 1},
        {"text": "Cost management", "level": 1}
    ])
    print("✓ Slide 5: Why Change")
    
    # Slide 6: Why - Business Impact
    add_bullet_slide(prs, CONTENT_LAYOUT, "Why: Business Impact", [
        "Quantifiable benefits:",
        "",
        "Development Velocity",
        {"text": "90% less boilerplate code", "level": 1},
        {"text": "10x faster iterations", "level": 1},
        "",
        "Team Productivity",
        {"text": "New members productive in minutes", "level": 1},
        {"text": "Focus on models, not config", "level": 1},
        "",
        "Risk Reduction",
        {"text": "Fewer configuration errors", "level": 1},
        {"text": "Consistent security policies", "level": 1}
    ])
    print("✓ Slide 6: Business Impact")
    
    # Slide 7: How - Solution
    add_bullet_slide(prs, CONTENT_LAYOUT, "How: ML Platform SDK (mlp_sdk)", [
        "Configuration-driven SageMaker wrapper",
        "",
        "Core Concept",
        {"text": "Define infrastructure once in YAML", "level": 1},
        {"text": "Simple API for ML operations", "level": 1},
        {"text": "Override any setting at runtime", "level": 1},
        "",
        "Key Features",
        {"text": "Session-based with smart defaults", "level": 1},
        {"text": "Built-in governance & audit trails", "level": 1},
        {"text": "AES-256-GCM encryption", "level": 1},
        {"text": "Zero lock-in to underlying SDK", "level": 1}
    ])
    print("✓ Slide 7: Solution")
    
    # Slide 8: How - Architecture
    add_bullet_slide(prs, CONTENT_LAYOUT, "How: Three-Layer Architecture", [
        "Intelligent wrapper design:",
        "",
        "Layer 1: MLP_Session",
        {"text": "Configuration manager", "level": 1},
        {"text": "Session orchestrator", "level": 1},
        "",
        "Layer 2: Specialized Wrappers",
        {"text": "Training, Processing, Feature Store", "level": 1},
        {"text": "Pipeline orchestration", "level": 1},
        "",
        "Layer 3: SageMaker SDK v3",
        {"text": "Full access to underlying SDK", "level": 1},
        {"text": "No functionality locked away", "level": 1}
    ])
    print("✓ Slide 8: Architecture")
    
    # Slide 9: Configuration
    add_bullet_slide(prs, CONTENT_LAYOUT, "Configuration Precedence", [
        "Three-tier flexibility:",
        "",
        "Priority 1: Runtime Parameters",
        {"text": "Explicitly passed to methods", "level": 1},
        {"text": "Override for special cases", "level": 1},
        "",
        "Priority 2: YAML Configuration",
        {"text": "Team-wide defaults", "level": 1},
        {"text": "Environment-specific settings", "level": 1},
        "",
        "Priority 3: SDK Defaults",
        {"text": "Built-in SageMaker defaults", "level": 1},
        {"text": "Fallback when nothing specified", "level": 1}
    ])
    print("✓ Slide 9: Configuration")
    
    # Slide 10: Config File
    add_bullet_slide(prs, CONTENT_LAYOUT, "Configuration File Structure", [
        "Define once in admin-config.yaml:",
        "",
        {"text": "S3: buckets, prefixes", "level": 0},
        {"text": "Networking: VPC, subnets, security groups", "level": 0},
        {"text": "Compute: instance types, counts", "level": 0},
        {"text": "IAM: execution role ARN", "level": 0},
        {"text": "KMS: encryption keys (optional)", "level": 0},
        "",
        "Generate interactively:",
        {"text": "python generate_admin_config.py --interactive", "level": 1}
    ])
    print("✓ Slide 10: Config File")
    
    # Slide 11: Code Comparison
    add_two_column_slide(prs, TWO_COL_LAYOUT, "Before vs After: Training Job",
        [
            "❌ Without mlp_sdk",
            "",
            "50+ lines of code:",
            "• Import multiple classes",
            "• Define Compute config",
            "• Define Networking config",
            "• Define OutputDataConfig",
            "• Define ResourceConfig",
            "• Create ModelTrainer",
            "• Configure all parameters",
            "",
            "Verbose and error-prone"
        ],
        [
            "✅ With mlp_sdk",
            "",
            "5 lines of code:",
            "",
            "trainer = session.run_training_job(",
            "  job_name='my-job',",
            "  training_image=container,",
            "  inputs={'train': 's3://...'}",
            ")",
            "",
            "",
            "Clean and maintainable"
        ]
    )
    print("✓ Slide 11: Code Comparison")
    
    # Slide 12: Training Jobs
    add_bullet_slide(prs, CONTENT_LAYOUT, "Core Feature: Training Jobs", [
        "Simplified training execution:",
        "",
        "Basic Usage",
        {"text": "session = MLP_Session()", "level": 1},
        {"text": "trainer = session.run_training_job(...)", "level": 1},
        "",
        "Auto-Configured",
        {"text": "✓ Instance type from config", "level": 1},
        {"text": "✓ VPC and security setup", "level": 1},
        {"text": "✓ IAM role and KMS encryption", "level": 1},
        "",
        "Override when needed",
        {"text": "Use GPU for specific jobs", "level": 1}
    ])
    print("✓ Slide 12: Training")
    
    # Slide 13: Other Features
    add_bullet_slide(prs, CONTENT_LAYOUT, "More Features", [
        "Processing Jobs",
        {"text": "session.run_processing_job()", "level": 1},
        {"text": "Simplified data mapping", "level": 1},
        "",
        "Feature Store",
        {"text": "session.create_feature_group()", "level": 1},
        {"text": "Validates definitions", "level": 1},
        "",
        "Model Deployment",
        {"text": "session.deploy_model()", "level": 1},
        {"text": "Auto endpoint configuration", "level": 1}
    ])
    print("✓ Slide 13: More Features")
    
    # Slide 14: Advanced
    add_bullet_slide(prs, CONTENT_LAYOUT, "Advanced Capabilities", [
        "Audit Trails",
        {"text": "Track all operations", "level": 1},
        {"text": "Export to JSON or CSV", "level": 1},
        "",
        "Encryption",
        {"text": "AES-256-GCM for sensitive values", "level": 1},
        {"text": "AWS KMS integration", "level": 1},
        "",
        "Multi-Environment",
        {"text": "Separate configs for dev/staging/prod", "level": 1},
        "",
        "Full SDK Access",
        {"text": "No lock-in, use low-level APIs anytime", "level": 1}
    ])
    print("✓ Slide 14: Advanced")
    
    # Slide 15: Getting Started
    add_bullet_slide(prs, CONTENT_LAYOUT, "Getting Started (5 Minutes)", [
        "Three simple steps:",
        "",
        "Step 1: Install",
        {"text": "pip install mlp_sdk", "level": 1},
        "",
        "Step 2: Generate Config",
        {"text": "python generate_admin_config.py --interactive", "level": 1},
        "",
        "Step 3: Start Building",
        {"text": "from mlp_sdk import MLP_Session", "level": 1},
        {"text": "session = MLP_Session()", "level": 1},
        {"text": "trainer = session.run_training_job(...)", "level": 1}
    ])
    print("✓ Slide 15: Getting Started")
    
    # Save
    output_file = "ML_Platform_SDK_Deck.pptx"
    prs.save(output_file)
    
    print(f"\n{'='*60}")
    print(f"✅ Presentation created: {output_file}")
    print(f"   Total slides: {len(prs.slides)}")
    print(f"   Bullet points optimized for visibility")
    print(f"{'='*60}")
    
    return output_file

if __name__ == "__main__":
    try:
        output_file = create_presentation()
        print(f"\n🎉 Success! Open {output_file} to view.")
    except Exception as e:
        print(f"\n❌ Error: {e}")
        import traceback
        traceback.print_exc()
