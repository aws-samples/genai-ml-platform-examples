#!/usr/bin/env python3
"""
Generate PowerPoint presentation for ML Platform SDK
Uses AWS PowerPoint template and content from mlp_sdk documentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

# AWS Brand Colors
AWS_ORANGE = RGBColor(255, 153, 0)  # #FF9900
AWS_DARK = RGBColor(35, 47, 62)     # #232F3E
AWS_LIGHT_GRAY = RGBColor(234, 237, 237)  # #EAEDED
AWS_BLUE = RGBColor(0, 161, 222)    # #00A1DE

def load_template(template_path):
    """Load the AWS PowerPoint template"""
    if os.path.exists(template_path):
        return Presentation(template_path)
    else:
        print(f"Warning: Template not found at {template_path}, using blank presentation")
        return Presentation()

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    return slide

def add_section_header(prs, title):
    """Add section header slide"""
    slide_layout = prs.slide_layouts[2]  # Section header layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    return slide

def add_content_slide(prs, title, content_items):
    """Add content slide with bullet points"""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Add content
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    for item in content_items:
        if isinstance(item, dict):
            # Nested bullet
            p = tf.add_paragraph()
            p.text = item['text']
            p.level = item.get('level', 0)
            p.font.size = Pt(item.get('font_size', 18))
        else:
            # Simple bullet
            p = tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(18)
    
    return slide

def add_code_slide(prs, title, code_text, description=None):
    """Add slide with code example"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Add description if provided
    if description:
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(0.5)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = description
        p.font.size = Pt(14)
        p.font.color.rgb = AWS_DARK
    
    # Add code box
    left = Inches(0.5)
    top = Inches(2.2) if description else Inches(1.5)
    width = Inches(9)
    height = Inches(4.5) if description else Inches(5.2)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = code_text
    p.font.name = 'Courier New'
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(0, 0, 0)
    
    # Add background to code box
    fill = textbox.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)
    
    return slide

def add_comparison_slide(prs, title, before_title, before_items, after_title, after_items):
    """Add before/after comparison slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Before column
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(4.25)
    height = Inches(0.5)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = before_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = AWS_ORANGE
    
    # Before content
    top = Inches(2.4)
    height = Inches(4)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.clear()
    
    for item in before_items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.space_before = Pt(6)
    
    # After column
    left = Inches(5.25)
    top = Inches(1.8)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = after_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)
    
    # After content
    top = Inches(2.4)
    height = Inches(4)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.clear()
    
    for item in after_items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.space_before = Pt(6)
    
    return slide

def create_presentation():
    """Create the complete presentation"""
    
    # Load template
    template_path = "2026_aws_powerpoint_template_v1_07162eba.pptx"
    prs = load_template(template_path)
    
    print("Creating ML Platform SDK presentation...")
    
    # Slide 1: Title
    add_title_slide(
        prs,
        "Building an ML Platform SDK Wrapper",
        "Simplifying Machine Learning Operations with Configuration-Driven Infrastructure"
    )
    print("✓ Added title slide")
    
    # Slide 2: The Challenge
    add_section_header(prs, "The Challenge")
    
    add_content_slide(
        prs,
        "ML Infrastructure Complexity",
        [
            "Data scientists spend too much time on infrastructure",
            {"text": "VPC configurations, security groups, IAM roles", "level": 1},
            {"text": "S3 bucket policies, KMS encryption, network isolation", "level": 1},
            {"text": "Repeated boilerplate code across every project", "level": 1},
            "",
            "Inconsistency across teams and projects",
            {"text": "Different naming conventions", "level": 1},
            {"text": "Configuration drift between environments", "level": 1},
            {"text": "Knowledge silos and tribal knowledge", "level": 1},
            "",
            "Governance and compliance challenges",
            {"text": "Difficult to enforce security policies", "level": 1},
            {"text": "Limited audit trails", "level": 1},
            {"text": "Cost attribution problems", "level": 1}
        ]
    )
    print("✓ Added challenge slides")
    
    # Slide 3: The Solution
    add_section_header(prs, "The Solution: ML Platform SDK")
    
    add_content_slide(
        prs,
        "What is mlp_sdk?",
        [
            "Configuration-driven wrapper around SageMaker Python SDK v3",
            "",
            "Key Benefits:",
            {"text": "Define infrastructure once in YAML configuration", "level": 1},
            {"text": "Simple, declarative API for ML operations", "level": 1},
            {"text": "Runtime flexibility to override any setting", "level": 1},
            {"text": "Built-in governance with audit trails", "level": 1},
            {"text": "Zero lock-in - access underlying SDK anytime", "level": 1},
            "",
            "Result: 90% less boilerplate, 10x faster iteration"
        ]
    )
    print("✓ Added solution slides")
    
    # Slide 4: Architecture
    add_content_slide(
        prs,
        "Architecture Overview",
        [
            "Three-Layer Architecture:",
            "",
            {"text": "MLP_Session - Central orchestrator", "level": 1},
            {"text": "Configuration Manager - YAML config handling", "level": 1},
            {"text": "Specialized Wrappers - Training, Processing, Feature Store", "level": 1},
            "",
            "Configuration Precedence:",
            {"text": "1. Runtime Parameters (highest priority)", "level": 1},
            {"text": "2. YAML Configuration (middle priority)", "level": 1},
            {"text": "3. SageMaker SDK Defaults (lowest priority)", "level": 1}
        ]
    )
    print("✓ Added architecture slide")
    
    # Slide 5: Configuration Example
    add_code_slide(
        prs,
        "Configuration File (admin-config.yaml)",
        """defaults:
  # S3 Configuration
  s3:
    default_bucket: "my-ml-platform-bucket"
    input_prefix: "data/input/"
    output_prefix: "data/output/"
    
  # Networking Configuration  
  networking:
    vpc_id: "vpc-0a1b2c3d"
    security_group_ids: ["sg-12345678"]
    subnets: ["subnet-12345678", "subnet-87654321"]
    
  # Compute Configuration
  compute:
    processing_instance_type: "ml.m5.large"
    training_instance_type: "ml.m5.xlarge"
    
  # IAM Configuration
  iam:
    execution_role: "arn:aws:iam::123456789012:role/SageMakerRole"
    
  # KMS Configuration (optional)
  kms:
    key_id: "arn:aws:kms:REGION:ACCOUNT-ID:key/KEY-ID" """,
        "Define infrastructure once, use everywhere"
    )
    print("✓ Added configuration slide")
    
    # Slide 6: Before/After Comparison
    add_comparison_slide(
        prs,
        "Code Comparison: Training Job",
        "❌ Without mlp_sdk (50+ lines)",
        [
            "• Import multiple SageMaker classes",
            "• Define Compute configuration",
            "• Define Networking configuration",
            "• Define OutputDataConfig",
            "• Define ResourceConfig",
            "• Define StoppingCondition",
            "• Create ModelTrainer",
            "• Configure all parameters",
            "• Start training",
            "",
            "Result: Verbose, error-prone, repetitive"
        ],
        "✅ With mlp_sdk (5 lines)",
        [
            "trainer = session.run_training_job(",
            "    job_name='my-training',",
            "    training_image=container,",
            "    inputs={'train': 's3://bucket/data/'}",
            ")",
            "",
            "",
            "",
            "",
            "Result: Clean, focused, maintainable"
        ]
    )
    print("✓ Added comparison slide")
    
    # Slide 7: Basic Usage
    add_code_slide(
        prs,
        "Basic Usage: Getting Started",
        """# Install
pip install mlp_sdk

# Initialize session (loads configuration automatically)
from mlp_sdk import MLP_Session
session = MLP_Session()

# Run training job - infrastructure from config!
trainer = session.run_training_job(
    job_name="xgboost-training",
    training_image=container_uri,
    inputs={
        'train': 's3://my-bucket/train/',
        'validation': 's3://my-bucket/validation/'
    },
    hyperparameters={
        'objective': 'binary:logistic',
        'num_round': '100',
        'max_depth': '5'
    }
)

# Instance type, VPC, role, encryption all from config!""",
        "Simple, declarative API"
    )
    print("✓ Added basic usage slide")
    
    # Slide 8: Training Example
    add_code_slide(
        prs,
        "Example: XGBoost Training",
        """# Start training with minimal code
trainer = session.run_training_job(
    job_name=f"xgboost-training-{timestamp}",
    training_image=xgboost_container,
    inputs={
        'train': train_s3_path,
        'validation': val_s3_path
    },
    hyperparameters={
        'objective': 'binary:logistic',
        'num_round': '100',
        'max_depth': '5',
        'eta': '0.2'
    }
)

# What mlp_sdk handled automatically:
# ✅ Instance type (from config)
# ✅ VPC configuration (from config)
# ✅ Security groups (from config)
# ✅ IAM role (from config)
# ✅ KMS encryption (from config)""",
        "Focus on ML, not infrastructure"
    )
    print("✓ Added training example slide")
    
    # Slide 9: Processing Jobs
    add_code_slide(
        prs,
        "Processing Jobs",
        """# Run data preprocessing
processor = session.run_processing_job(
    job_name="data-preprocessing",
    processing_script="scripts/preprocess.py",
    inputs=[{
        "source": "s3://my-bucket/raw-data/",
        "destination": "/opt/ml/processing/input"
    }],
    outputs=[{
        "source": "/opt/ml/processing/output",
        "destination": "s3://my-bucket/processed-data/"
    }]
)

# Override instance type for large datasets
processor = session.run_processing_job(
    job_name="large-preprocessing",
    processing_script="scripts/preprocess.py",
    inputs=[...],
    outputs=[...],
    instance_type="ml.m5.4xlarge"  # Runtime override
)""",
        "Simple processing with runtime flexibility"
    )
    print("✓ Added processing slide")
    
    # Slide 10: Feature Store
    add_code_slide(
        prs,
        "Feature Store Operations",
        """# Create feature group with defaults
feature_group = session.create_feature_group(
    feature_group_name="customer-features",
    record_identifier_name="customer_id",
    event_time_feature_name="event_time",
    feature_definitions=[
        {"FeatureName": "customer_id", "FeatureType": "String"},
        {"FeatureName": "age", "FeatureType": "Integral"},
        {"FeatureName": "income", "FeatureType": "Fractional"},
        {"FeatureName": "event_time", "FeatureType": "String"}
    ]
)

# Enable online store for real-time features
feature_group = session.create_feature_group(
    feature_group_name="realtime-features",
    record_identifier_name="id",
    event_time_feature_name="timestamp",
    feature_definitions=[...],
    enable_online_store=True  # Override config
)""",
        "Simplified feature store management"
    )
    print("✓ Added feature store slide")
    
    # Slide 11: Advanced Features
    add_content_slide(
        prs,
        "Advanced Features",
        [
            "Audit Trails",
            {"text": "Track all operations for debugging and compliance", "level": 1},
            {"text": "Export to JSON or CSV for analysis", "level": 1},
            "",
            "Encryption Support",
            {"text": "AES-256-GCM encryption for sensitive config values", "level": 1},
            {"text": "Integration with AWS KMS", "level": 1},
            "",
            "Multi-Environment Support",
            {"text": "Separate configs for dev/staging/prod", "level": 1},
            {"text": "Environment-specific settings", "level": 1},
            "",
            "Access to Underlying SDK",
            {"text": "No lock-in - access SageMaker SDK anytime", "level": 1},
            {"text": "Use low-level APIs when needed", "level": 1}
        ]
    )
    print("✓ Added advanced features slide")
    
    # Slide 12: Audit Trails
    add_code_slide(
        prs,
        "Audit Trails: Complete Visibility",
        """# Enable audit trail (enabled by default)
session = MLP_Session(enable_audit_trail=True)

# Perform operations
session.run_training_job(...)
session.run_processing_job(...)
session.create_feature_group(...)

# Get audit trail
entries = session.get_audit_trail()
print(f"Total operations: {len(entries)}")

# Filter by operation type
training_ops = session.get_audit_trail(operation="run_training_job")

# Filter by status
failed_ops = session.get_audit_trail(status="failed")

# Export for analysis
session.export_audit_trail("audit-trail.json", format="json")
session.export_audit_trail("audit-trail.csv", format="csv")""",
        "Track operations for debugging, compliance, and cost attribution"
    )
    print("✓ Added audit trails slide")
    
    # Slide 13: Encryption
    add_code_slide(
        prs,
        "Encryption: Secure Your Configuration",
        """from mlp_sdk.config import ConfigurationManager

# Generate encryption key
key = ConfigurationManager.generate_key()

# Encrypt sensitive fields
config_manager = ConfigurationManager(encryption_key=key)
config_manager.encrypt_config_file(
    input_path="config.yaml",
    output_path="config-encrypted.yaml",
    fields_to_encrypt=[
        "defaults.iam.execution_role",
        "defaults.kms.key_id"
    ]
)

# Load encrypted configuration
session = MLP_Session(
    config_path="config-encrypted.yaml",
    encryption_key=key
)

# Key management options:
# - Environment variables
# - File-based storage
# - AWS KMS integration""",
        "Protect sensitive configuration values"
    )
    print("✓ Added encryption slide")
    
    # Slide 14: Benefits
    add_content_slide(
        prs,
        "Key Benefits",
        [
            "Faster Development Cycles",
            {"text": "90% less boilerplate code", "level": 1},
            {"text": "10x faster iteration", "level": 1},
            "",
            "Consistent Team Standards",
            {"text": "Shared configuration across teams", "level": 1},
            {"text": "Enforced naming conventions and security policies", "level": 1},
            "",
            "Reduced Errors",
            {"text": "Type validation catches errors early", "level": 1},
            {"text": "Clear, actionable error messages", "level": 1},
            "",
            "Seamless Onboarding",
            {"text": "New team members productive in minutes", "level": 1},
            {"text": "No need to learn infrastructure details", "level": 1},
            "",
            "Enterprise-Ready",
            {"text": "Encryption, audit trails, multi-environment support", "level": 1}
        ]
    )
    print("✓ Added benefits slide")
    
    # Slide 15: Use Cases
    add_content_slide(
        prs,
        "Real-World Use Cases",
        [
            "Model Training Pipelines",
            {"text": "Standardize training across teams", "level": 1},
            {"text": "Consistent hyperparameter tracking", "level": 1},
            "",
            "Data Processing Workflows",
            {"text": "Unified preprocessing infrastructure", "level": 1},
            {"text": "Reusable processing scripts", "level": 1},
            "",
            "Feature Engineering",
            {"text": "Centralized feature store management", "level": 1},
            {"text": "Consistent feature definitions", "level": 1},
            "",
            "Multi-Environment Deployments",
            {"text": "Dev/staging/prod with different configs", "level": 1},
            {"text": "Promote models across environments", "level": 1},
            "",
            "Team Collaboration",
            {"text": "Shared infrastructure standards", "level": 1},
            {"text": "Reduced knowledge silos", "level": 1}
        ]
    )
    print("✓ Added use cases slide")
    
    # Slide 16: Getting Started
    add_code_slide(
        prs,
        "Getting Started (5 Minutes)",
        """# 1. Install mlp_sdk
pip install mlp_sdk

# 2. Generate configuration
python examples/generate_admin_config.py --interactive

# 3. Initialize session
from mlp_sdk import MLP_Session
session = MLP_Session()

# 4. Start building!
trainer = session.run_training_job(
    job_name="my-first-training",
    training_image=container_uri,
    inputs={'train': 's3://bucket/data/'}
)

# That's it! You're ready to go.

# Explore examples:
# - examples/basic_usage.py
# - examples/xgboost_training_example.ipynb
# - examples/sagemaker_operations.py""",
        "From zero to production in minutes"
    )
    print("✓ Added getting started slide")
    
    # Slide 17: Best Practices
    add_content_slide(
        prs,
        "Best Practices",
        [
            "Configuration Management",
            {"text": "Use version control for configuration files", "level": 1},
            {"text": "Separate configs for each environment", "level": 1},
            {"text": "Encrypt sensitive values", "level": 1},
            "",
            "Team Collaboration",
            {"text": "Share configuration across team", "level": 1},
            {"text": "Document custom settings", "level": 1},
            {"text": "Use consistent naming conventions", "level": 1},
            "",
            "Security",
            {"text": "Enable audit trails for compliance", "level": 1},
            {"text": "Use IAM roles with least privilege", "level": 1},
            {"text": "Enable encryption for sensitive data", "level": 1},
            "",
            "Development Workflow",
            {"text": "Start with defaults, override when needed", "level": 1},
            {"text": "Test in dev before promoting to prod", "level": 1}
        ]
    )
    print("✓ Added best practices slide")
    
    # Slide 18: Summary
    add_content_slide(
        prs,
        "Summary",
        [
            "mlp_sdk transforms ML operations:",
            "",
            {"text": "✅ 90% less boilerplate code", "level": 1},
            {"text": "✅ Configuration-driven infrastructure", "level": 1},
            {"text": "✅ Runtime flexibility when needed", "level": 1},
            {"text": "✅ Built-in governance and audit trails", "level": 1},
            {"text": "✅ Zero lock-in to underlying SDK", "level": 1},
            "",
            "Focus on building better models, not managing infrastructure",
            "",
            "Get started today:",
            {"text": "pip install mlp_sdk", "level": 1},
            {"text": "GitHub: github.com/example/mlp_sdk", "level": 1},
            {"text": "Documentation: mlp-sdk.readthedocs.io", "level": 1}
        ]
    )
    print("✓ Added summary slide")
    
    # Slide 19: Q&A
    add_section_header(prs, "Questions?")
    print("✓ Added Q&A slide")
    
    # Save presentation
    output_file = "ML_Platform_SDK_Presentation.pptx"
    prs.save(output_file)
    print(f"\n✅ Presentation created: {output_file}")
    print(f"   Total slides: {len(prs.slides)}")
    
    return output_file

if __name__ == "__main__":
    try:
        output_file = create_presentation()
        print(f"\n🎉 Success! Open {output_file} to view the presentation.")
    except Exception as e:
        print(f"\n❌ Error creating presentation: {e}")
        import traceback
        traceback.print_exc()
