#!/usr/bin/env python3
"""
Analyze updated template and create ML Platform SDK presentation
with correct slide layouts
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def analyze_template(template_path):
    """Analyze template to understand available layouts"""
    prs = Presentation(template_path)
    
    print(f"\n{'='*70}")
    print(f"TEMPLATE ANALYSIS: {template_path}")
    print(f"{'='*70}\n")
    
    print(f"Total slide layouts: {len(prs.slide_layouts)}\n")
    
    # Categorize layouts
    title_layouts = []
    agenda_layouts = []
    content_layouts = []
    two_column_layouts = []
    section_layouts = []
    
    for idx, layout in enumerate(prs.slide_layouts):
        name_lower = layout.name.lower()
        print(f"Layout {idx}: {layout.name}")
        print(f"  Placeholders: {len(layout.placeholders)}")
        
        # Show placeholder details
        for placeholder in layout.placeholders:
            ph_type = placeholder.placeholder_format.type
            print(f"    [{placeholder.placeholder_format.idx}] {placeholder.name} (type: {ph_type})")
        
        # Categorize
        if 'title' in name_lower and 'slide' in name_lower:
            title_layouts.append((idx, layout.name))
        elif 'agenda' in name_lower:
            agenda_layouts.append((idx, layout.name))
        elif 'two' in name_lower or 'comparison' in name_lower:
            two_column_layouts.append((idx, layout.name))
        elif 'section' in name_lower or 'header' in name_lower:
            section_layouts.append((idx, layout.name))
        elif 'content' in name_lower or 'bullet' in name_lower:
            content_layouts.append((idx, layout.name))
        
        print()
    
    # Print categorization
    print(f"\n{'='*70}")
    print("LAYOUT CATEGORIZATION")
    print(f"{'='*70}\n")
    
    print(f"Title Layouts ({len(title_layouts)}):")
    for idx, name in title_layouts:
        print(f"  [{idx}] {name}")
    
    print(f"\nAgenda Layouts ({len(agenda_layouts)}):")
    for idx, name in agenda_layouts:
        print(f"  [{idx}] {name}")
    
    print(f"\nContent Layouts ({len(content_layouts)}):")
    for idx, name in content_layouts:
        print(f"  [{idx}] {name}")
    
    print(f"\nTwo Column Layouts ({len(two_column_layouts)}):")
    for idx, name in two_column_layouts:
        print(f"  [{idx}] {name}")
    
    print(f"\nSection Header Layouts ({len(section_layouts)}):")
    for idx, name in section_layouts:
        print(f"  [{idx}] {name}")
    
    print(f"\n{'='*70}\n")
    
    return prs, {
        'title': title_layouts,
        'agenda': agenda_layouts,
        'content': content_layouts,
        'two_column': two_column_layouts,
        'section': section_layouts
    }

def add_slide_with_layout(prs, layout_idx, title_text, content=None):
    """Add slide using specific layout index"""
    slide_layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    
    # Add content if provided
    if content:
        # Find content placeholder - try different types
        content_added = False
        
        # First try to find a shape with text_frame
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                tf = shape.text_frame
                tf.clear()
                
                for item in content:
                    if isinstance(item, dict):
                        p = tf.add_paragraph()
                        p.text = item['text']
                        p.level = item.get('level', 0)
                        if p.font:
                            p.font.size = Pt(item.get('font_size', 18))
                    elif item == "":
                        p = tf.add_paragraph()
                        p.text = ""
                    else:
                        p = tf.add_paragraph()
                        p.text = item
                        p.level = 0
                        if p.font:
                            p.font.size = Pt(18)
                
                content_added = True
                break
        
        # If no text_frame found, try placeholders
        if not content_added:
            for placeholder in slide.placeholders:
                ph_type = placeholder.placeholder_format.type
                # Try type 2 (BODY) or type 7 (OBJECT)
                if ph_type in [2, 7]:
                    try:
                        if placeholder.has_text_frame:
                            tf = placeholder.text_frame
                            tf.clear()
                            
                            for item in content:
                                if isinstance(item, dict):
                                    p = tf.add_paragraph()
                                    p.text = item['text']
                                    p.level = item.get('level', 0)
                                    if p.font:
                                        p.font.size = Pt(item.get('font_size', 18))
                                elif item == "":
                                    p = tf.add_paragraph()
                                    p.text = ""
                                else:
                                    p = tf.add_paragraph()
                                    p.text = item
                                    p.level = 0
                                    if p.font:
                                        p.font.size = Pt(18)
                            
                            content_added = True
                            break
                    except:
                        continue
        
        if not content_added:
            print(f"  Warning: Could not add content to slide '{title_text}'")
    
    return slide

def create_presentation_with_correct_layouts():
    """Create presentation using correct template layouts"""
    
    template_path = "ppt_template.pptx"
    
    # Analyze template
    print("Analyzing updated template...")
    prs, layout_categories = analyze_template(template_path)
    
    # Determine which layouts to use
    print("\n" + "="*70)
    print("CREATING PRESENTATION WITH CORRECT LAYOUTS")
    print("="*70 + "\n")
    
    # Select layouts
    title_layout_idx = layout_categories['title'][0][0] if layout_categories['title'] else 0
    agenda_layout_idx = layout_categories['agenda'][0][0] if layout_categories['agenda'] else 1
    # Use Layout 16 (Title and Bulleted Content) instead of Layout 15
    content_layout_idx = 16  # Title and Bulleted Content - better for bullet points
    two_col_layout_idx = layout_categories['two_column'][0][0] if layout_categories['two_column'] else content_layout_idx
    
    print(f"Using layouts:")
    print(f"  Title: Layout {title_layout_idx}")
    print(f"  Agenda: Layout {agenda_layout_idx}")
    print(f"  Content: Layout {content_layout_idx}")
    print(f"  Two Column: Layout {two_col_layout_idx}")
    print()
    
    # Slide 1: Title
    slide = add_slide_with_layout(prs, title_layout_idx, 
        "Building an ML Platform SDK Wrapper")
    # Add subtitle if placeholder exists
    for placeholder in slide.placeholders:
        if placeholder.placeholder_format.type == 3:  # Subtitle
            placeholder.text = "Simplifying SageMaker Operations with Configuration-Driven Infrastructure"
            break
    print("✓ Slide 1: Title")
    
    # Slide 2: Agenda
    add_slide_with_layout(prs, agenda_layout_idx, "Agenda", [
        "1. What: Understanding the Challenge",
        "2. Why: The Need for Simplification",
        "3. How: ML Platform SDK Architecture",
        "4. Configuration & Setup",
        "5. Core Features & Examples",
        "6. Advanced Capabilities",
        "7. Getting Started"
    ])
    print("✓ Slide 2: Agenda")
    
    # Slide 3: What - The Challenge
    add_slide_with_layout(prs, content_layout_idx, 
        "What: The ML Infrastructure Challenge", [
        "Data scientists face significant infrastructure overhead:",
        "",
        "Complex Configuration Requirements",
        {"text": "VPC settings, security groups, IAM roles", "level": 1},
        {"text": "S3 bucket policies, KMS encryption keys", "level": 1},
        {"text": "Network isolation and compliance settings", "level": 1},
        "",
        "Repetitive Boilerplate Code",
        {"text": "50+ lines of infrastructure code per training job", "level": 1},
        {"text": "Copy-paste errors across projects", "level": 1},
        {"text": "Inconsistent patterns across teams", "level": 1},
        "",
        "Governance Challenges",
        {"text": "Difficult to enforce security policies", "level": 1},
        {"text": "Limited audit trails and cost attribution", "level": 1}
    ])
    print("✓ Slide 3: The Challenge")
    
    # Slide 4: What - Current State
    add_slide_with_layout(prs, content_layout_idx,
        "What: Current State Without SDK Wrapper", [
        "Every SageMaker operation requires extensive setup:",
        "",
        "Training Job Example (Traditional Approach)",
        {"text": "Import 10+ SageMaker classes", "level": 1},
        {"text": "Define Compute, Networking, OutputDataConfig", "level": 1},
        {"text": "Configure ResourceConfig, StoppingCondition", "level": 1},
        {"text": "Specify VPC, subnets, security groups", "level": 1},
        {"text": "Set IAM roles and KMS encryption", "level": 1},
        {"text": "Create ModelTrainer with all parameters", "level": 1},
        "",
        "Result: 50+ lines of boilerplate",
        "Problem: Repeated across every project",
        "Impact: Slow iteration, high error rate"
    ])
    print("✓ Slide 4: Current State")
    
    # Slide 5: Why - The Need for Change
    add_slide_with_layout(prs, content_layout_idx,
        "Why: The Need for Simplification", [
        "Three Critical Problems to Solve:",
        "",
        "1. Infrastructure Abstraction Complexity",
        {"text": "Data scientists shouldn't manage VPCs and security groups", "level": 1},
        {"text": "Focus should be on ML, not infrastructure", "level": 1},
        "",
        "2. Consistency and Standardization",
        {"text": "Different teams use different patterns", "level": 1},
        {"text": "Configuration drift between environments", "level": 1},
        {"text": "Knowledge silos and tribal knowledge", "level": 1},
        "",
        "3. Governance and Compliance",
        {"text": "Central policy enforcement needed", "level": 1},
        {"text": "Audit trails for compliance", "level": 1},
        {"text": "Cost management and attribution", "level": 1}
    ])
    print("✓ Slide 5: Why Change")
    
    # Slide 6: Why - Business Impact
    add_slide_with_layout(prs, content_layout_idx,
        "Why: Business Impact", [
        "Quantifiable Benefits of SDK Wrapper:",
        "",
        "Development Velocity",
        {"text": "90% reduction in boilerplate code", "level": 1},
        {"text": "10x faster iteration cycles", "level": 1},
        {"text": "Hours to production instead of days", "level": 1},
        "",
        "Team Productivity",
        {"text": "New team members productive in minutes", "level": 1},
        {"text": "Reduced context switching", "level": 1},
        {"text": "Focus on model quality, not infrastructure", "level": 1},
        "",
        "Risk Reduction",
        {"text": "Fewer configuration errors", "level": 1},
        {"text": "Consistent security policies", "level": 1},
        {"text": "Complete audit trails", "level": 1}
    ])
    print("✓ Slide 6: Business Impact")
    
    # Slide 7: How - Solution Overview
    add_slide_with_layout(prs, content_layout_idx,
        "How: ML Platform SDK (mlp_sdk)", [
        "Configuration-driven wrapper around SageMaker Python SDK v3",
        "",
        "Core Concept",
        {"text": "Define infrastructure once in YAML configuration", "level": 1},
        {"text": "Simple, declarative API for ML operations", "level": 1},
        {"text": "Runtime flexibility to override any setting", "level": 1},
        "",
        "Key Features",
        {"text": "Session-based interface with smart defaults", "level": 1},
        {"text": "Built-in governance with audit trails", "level": 1},
        {"text": "AES-256-GCM encryption for sensitive values", "level": 1},
        {"text": "Zero lock-in - access underlying SDK anytime", "level": 1},
        "",
        "Result: 5 lines of code instead of 50+"
    ])
    print("✓ Slide 7: Solution Overview")
    
    # Slide 8: How - Architecture
    add_slide_with_layout(prs, content_layout_idx,
        "How: Three-Layer Architecture", [
        "Intelligent wrapper architecture:",
        "",
        "Layer 1: MLP_Session",
        {"text": "Central orchestrator and configuration manager", "level": 1},
        {"text": "Loads YAML configuration", "level": 1},
        {"text": "Manages boto3 and SageMaker sessions", "level": 1},
        "",
        "Layer 2: Specialized Wrappers",
        {"text": "Training Wrapper - Model training operations", "level": 1},
        {"text": "Processing Wrapper - Data preprocessing", "level": 1},
        {"text": "Feature Store Wrapper - Feature management", "level": 1},
        {"text": "Pipeline Wrapper - ML workflow orchestration", "level": 1},
        "",
        "Layer 3: SageMaker SDK v3",
        {"text": "Full access to underlying SDK", "level": 1},
        {"text": "No functionality locked away", "level": 1}
    ])
    print("✓ Slide 8: Architecture")
    
    # Slide 9: How - Configuration Precedence
    add_slide_with_layout(prs, content_layout_idx,
        "How: Configuration Precedence System", [
        "Three-tier precedence for maximum flexibility:",
        "",
        "Priority 1: Runtime Parameters (Highest)",
        {"text": "Explicitly passed to wrapper methods", "level": 1},
        {"text": "Override everything for special cases", "level": 1},
        {"text": "Example: instance_type='ml.p3.8xlarge'", "level": 1},
        "",
        "Priority 2: YAML Configuration (Middle)",
        {"text": "Defined in admin-config.yaml", "level": 1},
        {"text": "Team-wide defaults and standards", "level": 1},
        {"text": "Environment-specific settings", "level": 1},
        "",
        "Priority 3: SageMaker SDK Defaults (Lowest)",
        {"text": "Built-in SageMaker defaults", "level": 1},
        {"text": "Fallback when nothing specified", "level": 1}
    ])
    print("✓ Slide 9: Configuration Precedence")
    
    # Slide 10: Configuration Example
    add_slide_with_layout(prs, content_layout_idx,
        "Configuration File Structure", [
        "Define infrastructure once in admin-config.yaml:",
        "",
        "S3 Configuration",
        {"text": "default_bucket, input_prefix, output_prefix", "level": 1},
        "",
        "Networking Configuration",
        {"text": "vpc_id, security_group_ids, subnets", "level": 1},
        "",
        "Compute Configuration",
        {"text": "processing_instance_type, training_instance_type", "level": 1},
        {"text": "instance_count, volume_size", "level": 1},
        "",
        "IAM Configuration",
        {"text": "execution_role ARN", "level": 1},
        "",
        "KMS Configuration (Optional)",
        {"text": "key_id for encryption", "level": 1}
    ])
    print("✓ Slide 10: Configuration")
    
    # Slide 11: Code Comparison (Two Column)
    slide = add_slide_with_layout(prs, two_col_layout_idx, 
        "Before vs After: Training Job")
    
    # Add content to two columns
    # Find all content placeholders (type 7 = OBJECT)
    content_placeholders = [p for p in slide.placeholders if p.placeholder_format.type == 7]
    
    if len(content_placeholders) >= 2:
        # Left column
        if content_placeholders[0].has_text_frame:
            tf = content_placeholders[0].text_frame
            tf.clear()
            left_content = [
                "❌ Without mlp_sdk (50+ lines)",
                "",
                "• Import multiple classes",
                "• Define Compute config",
                "• Define Networking config",
                "• Define OutputDataConfig",
                "• Define ResourceConfig",
                "• Define StoppingCondition",
                "• Create ModelTrainer",
                "• Configure all parameters",
                "• Start training",
                "",
                "Verbose, error-prone, repetitive"
            ]
            for item in left_content:
                p = tf.add_paragraph()
                p.text = item
                if p.font:
                    p.font.size = Pt(16)
        
        # Right column
        if content_placeholders[1].has_text_frame:
            tf = content_placeholders[1].text_frame
            tf.clear()
            right_content = [
                "✅ With mlp_sdk (5 lines)",
                "",
                "trainer = session.run_training_job(",
                "    job_name='my-training',",
                "    training_image=container,",
                "    inputs={'train': 's3://...'}",
                ")",
                "",
                "",
                "",
                "",
                "Clean, focused, maintainable"
            ]
            for item in right_content:
                p = tf.add_paragraph()
                p.text = item
                if p.font:
                    p.font.size = Pt(16)
    
    print("✓ Slide 11: Code Comparison")
    
    # Slide 12: Core Features - Training
    add_slide_with_layout(prs, content_layout_idx,
        "Core Feature: Training Jobs", [
        "Simplified training job execution:",
        "",
        "Basic Usage",
        {"text": "session = MLP_Session()", "level": 1},
        {"text": "trainer = session.run_training_job(...)", "level": 1},
        "",
        "What's Handled Automatically",
        {"text": "✓ Instance type from configuration", "level": 1},
        {"text": "✓ VPC and security group setup", "level": 1},
        {"text": "✓ IAM role assignment", "level": 1},
        {"text": "✓ KMS encryption configuration", "level": 1},
        {"text": "✓ Output path management", "level": 1},
        "",
        "Runtime Flexibility",
        {"text": "Override any setting when needed", "level": 1},
        {"text": "Example: Use GPU for specific job", "level": 1}
    ])
    print("✓ Slide 12: Training Jobs")
    
    # Slide 13: Core Features - Processing & Feature Store
    add_slide_with_layout(prs, content_layout_idx,
        "Core Features: Processing & Feature Store", [
        "Processing Jobs",
        {"text": "session.run_processing_job()", "level": 1},
        {"text": "Simplified input/output data mapping", "level": 1},
        {"text": "Automatic resource configuration", "level": 1},
        "",
        "Feature Store Operations",
        {"text": "session.create_feature_group()", "level": 1},
        {"text": "Validates feature definitions", "level": 1},
        {"text": "Handles online/offline store config", "level": 1},
        {"text": "Manages feature group metadata", "level": 1},
        "",
        "Model Deployment",
        {"text": "session.deploy_model()", "level": 1},
        {"text": "Automatic endpoint configuration", "level": 1},
        {"text": "VPC and security settings applied", "level": 1}
    ])
    print("✓ Slide 13: Processing & Feature Store")
    
    # Slide 14: Advanced Features
    add_slide_with_layout(prs, content_layout_idx,
        "Advanced Capabilities", [
        "Audit Trails",
        {"text": "Track all operations for debugging and compliance", "level": 1},
        {"text": "Export to JSON or CSV for analysis", "level": 1},
        {"text": "Filter by operation type or status", "level": 1},
        "",
        "Encryption Support",
        {"text": "AES-256-GCM for sensitive configuration values", "level": 1},
        {"text": "Integration with AWS KMS", "level": 1},
        {"text": "Environment variables or file-based keys", "level": 1},
        "",
        "Multi-Environment Support",
        {"text": "Separate configs for dev/staging/prod", "level": 1},
        {"text": "Environment-specific settings", "level": 1},
        "",
        "Full SDK Access",
        {"text": "No lock-in - access SageMaker SDK anytime", "level": 1},
        {"text": "Use low-level APIs when needed", "level": 1}
    ])
    print("✓ Slide 14: Advanced Features")
    
    # Slide 15: Getting Started
    add_slide_with_layout(prs, content_layout_idx,
        "Getting Started (5 Minutes)", [
        "Quick start in three steps:",
        "",
        "Step 1: Install",
        {"text": "pip install mlp_sdk", "level": 1},
        "",
        "Step 2: Generate Configuration",
        {"text": "python examples/generate_admin_config.py --interactive", "level": 1},
        {"text": "Answer prompts to create admin-config.yaml", "level": 1},
        "",
        "Step 3: Start Building",
        {"text": "from mlp_sdk import MLP_Session", "level": 1},
        {"text": "session = MLP_Session()", "level": 1},
        {"text": "trainer = session.run_training_job(...)", "level": 1},
        "",
        "Resources",
        {"text": "Examples: examples/xgboost_training_example.ipynb", "level": 1},
        {"text": "Documentation: README.md", "level": 1}
    ])
    print("✓ Slide 15: Getting Started")
    
    # Save presentation
    output_file = "ML_Platform_SDK_Deck.pptx"
    prs.save(output_file)
    
    print("="*70)
    print(f"\n✅ Presentation created: {output_file}")
    print(f"   Total slides: {len(prs.slides)}")
    print(f"   Template used: {template_path}")
    print(f"\n   Layouts used:")
    print(f"     Title: Layout {title_layout_idx}")
    print(f"     Agenda: Layout {agenda_layout_idx}")
    print(f"     Content: Layout {content_layout_idx}")
    print(f"     Two Column: Layout {two_col_layout_idx}")
    
    return output_file

if __name__ == "__main__":
    try:
        output_file = create_presentation_with_correct_layouts()
        print(f"\n🎉 Success! Open {output_file} to view the presentation.")
    except Exception as e:
        print(f"\n❌ Error creating presentation: {e}")
        import traceback
        traceback.print_exc()
